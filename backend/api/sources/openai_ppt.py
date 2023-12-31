import glob
import os
from pathlib import Path
import random
import re
import string
import asyncio
import httpx
import json
import uuid
from datetime import datetime, timezone
from typing import Optional

from icrawler import ImageDownloader
from icrawler.builtin import BingImageCrawler
from pptx import Presentation

from api.conf import Config, Credentials
from api.enums import OpenaiApiChatModels, ChatSourceTypes
from api.exceptions import OpenaiApiException
from api.models.doc import OpenaiApiChatMessage, OpenaiApiConversationHistoryDocument, OpenaiApiChatMessageMetadata, \
    OpenaiApiChatMessageTextContent
from api.schemas.openai_schemas import OpenaiChatResponse
from utils.common import singleton_with_lock
from utils.logger import get_logger

import openai

config = Config()
credentials = Credentials()

async def _check_response(response: httpx.Response) -> None:
    # 改成自带的错误处理
    try:
        response.raise_for_status()
    except httpx.HTTPStatusError as ex:
        await response.aread()
        error = OpenaiApiException(
            message=response.text,
            code=response.status_code,
        )
        raise error from ex


def make_session() -> httpx.AsyncClient:
    if config.openai_api.proxy is not None:
        proxies = {
            "http://": config.openai_api.proxy,
            "https://": config.openai_api.proxy,
        }
        session = httpx.AsyncClient(proxies=proxies, timeout=None)
    else:
        session = httpx.AsyncClient(timeout=None)
    return session


openai.api_key = credentials.openai_api_key

if config.openai_api.proxy is not None:
    openai.proxy = {'http': config.openai_api.proxy, 'https': config.openai_api.proxy}

bad_coding_practice = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits) for _ in
                              range(16))


def refresh_bad_coding_practice():
    global bad_coding_practice
    bad_coding_practice = ''.join(random.choice(string.ascii_uppercase + string.ascii_lowercase + string.digits)
                                  for _ in range(16))
    return


class PrefixNameDownloader(ImageDownloader):

    def get_filename(self, task, default_ext):
        filename = super(PrefixNameDownloader, self).get_filename(
            task, default_ext)
        print(bad_coding_practice)
        return 'prefix_' + bad_coding_practice + filename


async def generate_ppt(file_path, topic, slide_length, user_id: str):
    folder = Path() / "data" / "caches" / user_id
    folder.mkdir(parents=True, exist_ok=True)

    root = Presentation(file_path)

    message = f"""
    Create content for a slideshow presentation.
    The content's topic is {topic}. 
    The slideshow is {slide_length} slides long. 
    The content is written in the Chinese.


    You are allowed to use the following slide types:

    Slide types:
    Title Slide - (Title, Subtitle)
    Content Slide - (Title, Content)
    Image Slide - (Title, Content, Image)
    Thanks Slide - (Title)

    Put this tag before the Title Slide: [L_TS]
    Put this tag before the Content Slide: [L_CS]
    Put this tag before the Thanks Slide: [L_THS]

    Put "[SLIDEBREAK]" after each slide 

    For example:
    [L_TS]
    [TITLE]Mental Health[/TITLE]

    [SLIDEBREAK]

    [L_CS] 
    [TITLE]Mental Health Definition[/TITLE]
    [CONTENT]
    1. Definition: A person’s condition with regard to their psychological and emotional well-being
    2. Can impact one's physical health
    3. Stigmatized too often.
    [/CONTENT]

    [SLIDEBREAK]

    Put this tag before the Title: [TITLE]
    Put this tag after the Title: [/TITLE]
    Put this tag before the Subitle: [SUBTITLE]
    Put this tag after the Subtitle: [/SUBTITLE]
    Put this tag before the Content: [CONTENT]
    Put this tag after the Content: [/CONTENT]

    Elaborate on the Content, provide as much information as possible.
    You put a [/CONTENT] at the end of the Content.
    Do not reply as if you are talking about the slideshow itself. (ex. "Include pictures here about...")
    Do not include any special characters (?, !, ., :, ) in the Title.
    Do not include any additional information in your response and stick to the format."""

    response = await openai.ChatCompletion.acreate(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "user", "content": message}
        ]
    )

    def delete_all_slides():
        for i in range(len(root.slides) - 1, -1, -1):
            r_id = root.slides._sldIdLst[i].rId
            root.part.drop_rel(r_id)
            del root.slides._sldIdLst[i]

    def create_title_slide(title, subtitle):
        layout = root.slide_layouts[0]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def create_section_header_slide(title):
        layout = root.slide_layouts[2]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title

    def create_title_and_content_slide(title, content):
        layout = root.slide_layouts[1]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_title_and_content_and_image_slide(title, content, image_query):
        layout = root.slide_layouts[8]
        slide = root.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[2].text = content
        refresh_bad_coding_practice()
        baidu_crawler = BingImageCrawler(downloader_cls=PrefixNameDownloader, storage={'root_dir': folder})
        baidu_crawler.crawl(keyword=image_query, max_num=1)
        dir_path = os.path.dirname(os.path.realpath(__file__))
        file_name = glob.glob(f"prefix_{bad_coding_practice}*")
        print(file_name)
        img_path = os.path.join(dir_path, file_name[0])
        slide.shapes.add_picture(img_path, slide.placeholders[1].left, slide.placeholders[1].top,
                                 slide.placeholders[1].width, slide.placeholders[1].height)

    def find_text_in_between_tags(text, start_tag, end_tag):
        start_pos = text.find(start_tag)
        end_pos = text.find(end_tag)
        result = []
        while start_pos > -1 and end_pos > -1:
            text_between_tags = text[start_pos + len(start_tag):end_pos]
            result.append(text_between_tags)
            start_pos = text.find(start_tag, end_pos + len(end_tag))
            end_pos = text.find(end_tag, start_pos)
        res1 = "".join(result)
        res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]", '', res1)
        if len(result) > 0:
            return res2
        else:
            return ""

    def search_for_slide_type(text):
        tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
        found_text = next((s for s in tags if s in text), None)
        return found_text

    def parse_response(reply):
        list_of_slides = reply.split("[SLIDEBREAK]")
        for slide in list_of_slides:
            slide_type = search_for_slide_type(slide)
            if slide_type == "[L_TS]":
                create_title_slide(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]"),
                                   find_text_in_between_tags(str(slide), "[SUBTITLE]", "[/SUBTITLE]"))
            elif slide_type == "[L_CS]":

                create_title_and_content_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")),
                                               "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                 "[/CONTENT]")))

            elif slide_type == "[L_IS]":

                create_title_and_content_and_image_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]",
                                                                                           "[/TITLE]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[CONTENT]",
                                                                                           "[/CONTENT]")),
                                                         "".join(find_text_in_between_tags(str(slide), "[IMAGE]",
                                                                                           "[/IMAGE]")))
            elif slide_type == "[L_THS]":
                create_section_header_slide("".join(find_text_in_between_tags(str(slide), "[TITLE]", "[/TITLE]")))

    delete_all_slides()

    parse_response(response['choices'][0]['message']['content'])

    root.save(str(folder) + f"/{topic + user_id}.pptx")

    dir_path = str(folder)
    prefix = "prefix_"

    for file_name in os.listdir(dir_path):
        if file_name.startswith(prefix):
            file_path = os.path.join(dir_path, file_name)
            if os.path.isfile(file_path):
                os.remove(file_path)

    return os.getcwd() + "/" + str(folder) + f"/{topic + user_id}.pptx"



class OpenaiApiPptManager:    
    async def ask(content):
        if content == "" or content is None:
            return "内容不能为空"
    
        theme = content.split("：")[1].split("，")[0]
        topic = content.split("：")[2].split("，")[0]
        length = content.split("：")[3].split("，")[0]
        uuid="dsgfs"
        slides_limit=16
        filepath = './data/theme/' + theme
    
    
        if not os.path.exists(filepath):
            return f"PPT模版 {theme} 不存在，请重新输入！"
    
        if int(slides_limit) < int(length):
            return f"生成的PPT不能超过{slides_limit}页！"
    
        try:
            res = asyncio.run(generate_ppt(filepath, topic, length, uuid))
        except Exception as error:
            return str(error)
        url = "https://" + "stronge.yserver.top" + res

        return res+url
    
# 示例调用
content = "模版：1.pptx，主题：中国历史和中国人们为什么选择马克思主义，页数：16"
result = OpenaiApiPptManager.generate_ppt_sync(content)
print(result)
